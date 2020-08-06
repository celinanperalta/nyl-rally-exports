from pptx import Presentation
from pptx.util import Inches
import cv2
import os
from PIL import Image

def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size

    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width

    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)

    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # Placeholder width too wide:
    # if ratio_difference > 0:
    #     difference_on_each_side = ratio_difference / 2
    #     placeholder.crop_left = -difference_on_each_side
    #     placeholder.crop_right = -difference_on_each_side
    # # Placeholder height too high
    # else:
    #     difference_on_each_side = -ratio_difference / 2
    #     placeholder.crop_bottom = -difference_on_each_side
    #     placeholder.crop_top = -difference_on_each_side

def add_picture_slide(prs, path):
    layout8 = prs.slide_layouts[8]
    slide = prs.slides.add_slide(layout8)
    _add_image(slide, 1, path)

def add_title_slide(prs, title_text):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = ""

def add_bullet_slide(prs, title_text, subtitle, bullets):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title_text

    tf = body_shape.text_frame
    tf.text = subtitle

    n = len(bullets)

    for i in range(n):
        p = tf.add_paragraph()
        p.text = bullets[i]
        p.level = 1


prs = Presentation()

add_title_slide(prs, "RAID Slide Demo")
add_picture_slide(prs, "charts/owners.png")
add_bullet_slide(prs, "Bullet Slide", "Testing bullets!", ["a","b","c"])




prs.save('RAID.pptx')
